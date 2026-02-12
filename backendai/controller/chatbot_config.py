
import logging
from urllib.parse import urljoin
import PIL.Image
from bs4 import BeautifulSoup
import docx
from pptx import Presentation
import requests
from io import BytesIO
import PyPDF2
import base64
from typing import Optional
from services.openai_services import client
import boto3
import os

# Initialize S3 Client
s3_client = boto3.client(
    's3',
    aws_access_key_id=os.getenv('AWS_ACCESS_KEY_ID'),
    aws_secret_access_key=os.getenv('AWS_SECRET_ACCESS_KEY'),
    region_name=os.getenv('AWS_REGION')
)
BUCKET_NAME = os.getenv('S3_BUCKET_NAME')

def encode_image(image_bytes):
    return base64.b64encode(image_bytes).decode('utf-8')

def fetch_file_from_s3(key: str) -> Optional[bytes]:
    """Helper to fetch file content bytes from S3 securely"""
    try:
        response = s3_client.get_object(Bucket=BUCKET_NAME, Key=key)
        return response['Body'].read()
    except Exception as e:
        logging.error(f"Error fetching file from S3 (Key: {key}): {e}")
        return None

def generate_chatbot_training_instruction(content: str, image_bytes: Optional[bytes] = None) -> str:
    """
    Cleans and structures raw content (PDF, DOCX, HTML, PPT, etc.) for Vector DB Ingestion.
    Preserves 100% of information while fixing layout issues.
    """
    try:
        system_prompt = (
            "You are an expert Data Cleaner for AI RAG Systems. Your goal is to reformat raw extract "
            "into clean, structured Markdown for a Vector Database to index.\n"
            "CRITICAL INSTRUCTIONS:\n"
            "1. **PRESERVE ALL DATA**: Do NOT summarize, abbreviate, or omit any details. Keep every single fact, number, sentence, and paragraph.\n"
            "2. **Fix Scanned/PDF Text**: If the input has broken line breaks (typical in PDFs/Slides), merge them into coherent paragraphs. Fix disjointed sentences.\n"
            "3. **Structure It**: \n"
            "   - Use **Markdown Headers** (#, ##) for logical sections.\n"
            "   - Use **Bullet Points** for lists.\n"
            "   - Use **Markdown Tables** for any tabular data (preserve rows/cols accurately).\n"
            "   - Use **Code Blocks** for any code snippets.\n"
            "4. **Clean Layout**: Remove technical artifacts (like 'Page 1 of 10', 'Footer', 'Copyright', 'Menu', 'Nav'). Keep the BODY content 100% intact.\n"
            "5. **No Meta-Talk**: Do not add intros like 'Here is the cleaned data'. Output ONLY the cleaned content."
        )
         
        start_message = {
            "role": "user",
            "content": f"Here is the raw content:\n\n{content[:50000]}\n\nPlease reformat it cleanly while keeping ALL information."
        }
        
        if image_bytes:
            # OpenAI Vision Request
            base64_image = encode_image(image_bytes)
            start_message["content"] = [
                {"type": "text", "text": "Transcribe ALL text, tables, and details from this image accurately. Do not summarize. Maintain layout structure."},
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/jpeg;base64,{base64_image}"
                    }
                }
            ]

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": system_prompt},
                start_message
            ],
            temperature=0.3
        )
        
        return response.choices[0].message.content
    except Exception as e:
        logging.error(f"Error generating training instruction: {e}")
        return "Error generating summary."

def get_sitemap_urls(base_url: str) -> list[str]:
    sitemaps = ["/sitemap.xml", "/sitemap_index.xml"]
    
    for relative_path in sitemaps:
        try:
            sitemap_url = urljoin(base_url, relative_path)
            res = requests.get(sitemap_url, timeout=10)
            res.raise_for_status()

            soup = BeautifulSoup(res.text, "xml")
            urls = []

            # Normal sitemap
            for loc in soup.find_all("loc"):
                urls.append(loc.text.strip())
            
            if urls:
                print(f"Found sitemap at {sitemap_url} with {len(urls)} URLs")
                return urls

        except Exception as e:
            logging.warning(f"Could not read sitemap at {sitemap_url}: {e}")
            continue

    logging.error(f"Failed to find any sitemap for {base_url}")
    return []


def get_url_data(url):
    try:
        # Fetch data from the URL
        response = requests.get(url)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, 'html.parser')
        # Remove script and style elements
        for script_or_style in soup(['script', 'style', 'nav', 'footer']):
            script_or_style.decompose()

        # Get text
        text = soup.get_text(separator='\n')

        # Break into lines and remove leading and trailing space on each
        lines = (line.strip() for line in text.splitlines())
        # Break multi-headlines into a line each
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        # Drop blank lines
        url_data = '\n'.join(chunk for chunk in chunks if chunk)

        summary = generate_chatbot_training_instruction(url_data)
        print(summary)
        return summary

    except requests.exceptions.RequestException as e:
        return ""


# Fetch the image from S3 Key (formerly URL)
def image_data(s3_key: str):
    image_bytes = fetch_file_from_s3(s3_key)
    if image_bytes:
        summary = generate_chatbot_training_instruction("", image_bytes=image_bytes)
        return summary
    else:
        print(f"Failed to retrieve image: {s3_key}")

def pdf_data(s3_key: str):
    file_bytes = fetch_file_from_s3(s3_key)
    if file_bytes:
        file = BytesIO(file_bytes)
        reader = PyPDF2.PdfReader(file)
        text = ""
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]
            text += page.extract_text()
        
        summary = generate_chatbot_training_instruction(text)
        print(summary)
        return summary
        
        
def doc_data(s3_key: str):
    file_bytes = fetch_file_from_s3(s3_key)
    if file_bytes:
        file = BytesIO(file_bytes)
        doc = docx.Document(file)
        text = "\n".join([para.text for para in doc.paragraphs])
        
        summary = generate_chatbot_training_instruction(text)
        print(summary)
        return summary
        
def txt_data(s3_key: str):
    file_bytes = fetch_file_from_s3(s3_key)
    if file_bytes:
        text = file_bytes.decode('utf-8')
        
        summary = generate_chatbot_training_instruction(text)
        print(summary)
        return summary

def ppt_data(s3_key: str):
    file_bytes = fetch_file_from_s3(s3_key)
    if file_bytes:
        file = BytesIO(file_bytes)
        presentation = Presentation(file)
        text = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        summary = generate_chatbot_training_instruction(text)
        print(summary)
        return summary